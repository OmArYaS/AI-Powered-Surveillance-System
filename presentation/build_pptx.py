"""
Builds the .pptx presentation for the Sentinel graduation project.
Theme: Midnight Blue + Gold
All 25 slides with proper styling, animations notes, and speaker notes.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# Theme colors
BG = RGBColor(0x0B, 0x1A, 0x2E)         # Midnight blue
BG2 = RGBColor(0x0F, 0x22, 0x40)        # Lighter midnight
GOLD = RGBColor(0xD4, 0xAF, 0x37)       # Primary gold
GOLD_LIGHT = RGBColor(0xE8, 0xC8, 0x69) # Light gold
GOLD_DARK = RGBColor(0x9C, 0x7C, 0x1F)  # Dark gold
NAVY = RGBColor(0x1E, 0x3A, 0x5F)       # Medium navy
NAVY_LIGHT = RGBColor(0x2A, 0x4A, 0x75) # Light navy
TEXT = RGBColor(0xFF, 0xFF, 0xFF)       # White
TEXT_DIM = RGBColor(0xB8, 0xC5, 0xD3)   # Dim text
SUCCESS = RGBColor(0x4A, 0xDE, 0x80)    # Green
WARN = RGBColor(0xFB, 0xBF, 0x24)       # Amber
DANGER = RGBColor(0xEF, 0x44, 0x44)     # Red
INFO = RGBColor(0x60, 0xA5, 0xFA)       # Blue
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)     # Purple

ASSETS = os.path.join(os.path.dirname(__file__), "assets")
OUTPUT = os.path.join(os.path.dirname(__file__), "Sentinel_Presentation.pptx")

# Slide dimensions: 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=BG):
    """Set the background fill of a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, x, y, w, h, text, size=18, color=TEXT, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    """Add a text box with a single paragraph of text."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, bullets, size=18, color=TEXT,
                bullet_color=GOLD, line_spacing=1.3, font="Calibri"):
    """Add a text box with a list of bullets. Bullets can be (text, sub) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            main, sub = item
        else:
            main, sub = item, None

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing

        # Bullet character in gold
        bullet_run = p.add_run()
        bullet_run.text = "▸  "
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = bullet_color
        bullet_run.font.bold = True
        bullet_run.font.name = font

        # Main text
        main_run = p.add_run()
        main_run.text = main
        main_run.font.size = Pt(size)
        main_run.font.color.rgb = color
        main_run.font.name = font

        if sub:
            sp = tf.add_paragraph()
            sp.alignment = PP_ALIGN.LEFT
            sp.line_spacing = 1.0
            sub_run = sp.add_run()
            sub_run.text = "    " + sub
            sub_run.font.size = Pt(size - 4)
            sub_run.font.color.rgb = TEXT_DIM
            sub_run.font.italic = True
            sub_run.font.name = font
    return tb


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width if line_width else 1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rounded(slide, x, y, w, h, fill_color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_gold_accent_top(slide):
    """Top gold accent bar."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), GOLD)


def add_gold_accent_bottom(slide):
    """Bottom gold accent bar."""
    add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), GOLD)


def add_slide_title(slide, title, subtitle=None):
    """Standard slide title with gold underline."""
    add_text(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7),
             title, size=32, color=GOLD, bold=True)
    add_rect(slide, Inches(0.5), Inches(0.95), Inches(1.5), Inches(0.05), GOLD)
    if subtitle:
        add_text(slide, Inches(2.1), Inches(0.85), Inches(10.5), Inches(0.4),
                 subtitle, size=14, color=TEXT_DIM, italic=True)


def add_slide_footer(slide, page_num, total=25):
    """Slide number and project name in footer."""
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             "Sentinel · Smart Surveillance System", size=10, color=TEXT_DIM, italic=True)
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
             f"{page_num} / {total}", size=10, color=TEXT_DIM, italic=True, align=PP_ALIGN.RIGHT)


def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(path, x, y, width=w, height=h)
    elif w:
        return slide.shapes.add_picture(path, x, y, width=w)
    elif h:
        return slide.shapes.add_picture(path, x, y, height=h)
    return slide.shapes.add_picture(path, x, y)


def add_speaker_notes(slide, notes):
    """Add speaker notes to a slide."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes


# ======================================================================
# Build slides
# ======================================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ----------------------------------------------------------------
    # SLIDE 1 — Title
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s, BG)
    add_image(s, os.path.join(ASSETS, "title_bg.png"),
              Inches(0), Inches(0), w=SLIDE_W, h=SLIDE_H)

    # Decorative shield centered above title
    add_image(s, os.path.join(ASSETS, "logo.png"),
              Inches(5.665), Inches(1.0), w=Inches(2.0), h=Inches(2.0))

    # Main title
    add_text(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.9),
             "SENTINEL", size=72, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.6),
             "AI-Powered Smart Surveillance System", size=28, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.65), Inches(12.3), Inches(0.5),
             "Real-time Detection · Recognition · Threat & Theft Prevention",
             size=16, color=TEXT_DIM, italic=True, align=PP_ALIGN.CENTER)

    # Bottom info block
    add_rect(s, Inches(3.5), Inches(5.6), Inches(6.3), Inches(0.04), GOLD)

    add_text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.4),
             "Graduation Project · Defense Presentation", size=18, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
             "Team: [Your Name]    Supervisor: [Dr. Supervisor]",
             size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.4),
             "[Faculty · Department · 2025]",
             size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    add_speaker_notes(s, """Open strong. Greet the committee. Introduce yourself, your supervisor, and the project. Emphasize that Sentinel is not a toy demo — it's a full production-grade system that solves a real-world problem with measurable impact. Pause briefly after stating the name to let the title land.""")

    # ----------------------------------------------------------------
    # SLIDE 2 — Agenda
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Agenda", "What we will cover in the next 30 minutes")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 2)

    sections = [
        ("01", "Problem & Motivation", "Why surveillance needs AI"),
        ("02", "Project Objectives", "What we set out to build"),
        ("03", "System Architecture", "4-layer Clean Architecture"),
        ("04", "Core Features", "Detection · Recognition · Threat · Theft"),
        ("05", "Real-Time Pipeline", "End-to-end data flow"),
        ("06", "Performance & Metrics", "What we measured"),
        ("07", "Challenges & Solutions", "Problems we encountered"),
        ("08", "Impact & Use Cases", "Where this matters"),
        ("09", "Future Work", "What's next"),
        ("10", "Q & A", "Your questions"),
    ]
    for i, (num, title, sub) in enumerate(sections):
        y = Inches(1.5 + i * 0.52)
        # Number
        add_text(s, Inches(0.8), y, Inches(0.8), Inches(0.45),
                 num, size=18, color=GOLD, bold=True, font="Calibri")
        # Title
        add_text(s, Inches(1.7), y, Inches(5.0), Inches(0.45),
                 title, size=16, color=TEXT, bold=True)
        # Subtitle
        add_text(s, Inches(6.7), y, Inches(6.0), Inches(0.45),
                 sub, size=13, color=TEXT_DIM, italic=True)
        # Divider
        if i < len(sections) - 1:
            add_rect(s, Inches(1.7), y + Inches(0.5), Inches(11), Inches(0.01), NAVY)

    add_speaker_notes(s, """Walk the committee through what they will see. Highlight that you will go deep on each of the four core features. Set expectations for a 30-minute presentation followed by Q&A. The 'Challenges & Solutions' section is important — it shows engineering maturity, not just a happy-path demo.""")

    # ----------------------------------------------------------------
    # SLIDE 3 — Problem statement
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "The Problem", "Modern security threats outpace traditional surveillance")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 3)

    # Big quote
    add_text(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.0),
             '"Traditional CCTV is reactive — humans watch hundreds of feeds,\n and 95% of incidents are noticed only after they happen."',
             size=20, color=GOLD, italic=True, align=PP_ALIGN.CENTER)

    # 3 problem cards
    problems = [
        ("⚠", "Alert Fatigue", "Operators drown in 8+ hours of footage, miss critical events", DANGER, 2.5),
        ("⏱", "Slow Response", "Manual review takes 5-30 minutes per incident", WARN, 6.5),
        ("📉", "No Memory", "Cameras don't remember — the same person is anonymous every day", INFO, 10.5),
    ]
    for icon, title, desc, color, x in problems:
        add_rounded(s, Inches(x-1.2), Inches(3.3), Inches(3.6), Inches(2.5), BG2, color, 2)
        add_text(s, Inches(x-1.2), Inches(3.5), Inches(3.6), Inches(0.8),
                 icon, size=44, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(x-1.2), Inches(4.4), Inches(3.6), Inches(0.5),
                 title, size=18, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, Inches(x-1.2), Inches(4.9), Inches(3.6), Inches(0.8),
                 desc, size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Bottom call-out
    add_rect(s, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.5), NAVY)
    add_text(s, Inches(1.0), Inches(6.45), Inches(11.3), Inches(0.4),
             "We need intelligent, real-time, memory-aware surveillance that ALERTS, not RECORDS",
             size=14, color=GOLD, bold=True, italic=True, align=PP_ALIGN.CENTER)

    add_speaker_notes(s, """Set the stakes. The committee needs to understand why this work matters before you show the technical depth. Connect to real incidents: school shootings where nobody acted in time, retail theft costing $100B/year, lost children in malls. Make it human.""")

    # ----------------------------------------------------------------
    # SLIDE 4 — Project objectives
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Project Objectives", "Six concrete goals we delivered against")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 4)

    objectives = [
        ("1", "Real-time detection", "Identify people, phones, bags, weapons in <30 ms per frame", SUCCESS),
        ("2", "Face recognition", "Recognize registered staff; auto-register and re-identify unknowns", INFO),
        ("3", "Threat detection", "Fuse weapon, violence, and pose signals into severity-graded alerts", DANGER),
        ("4", "Theft detection", "Track personal items, infer ownership, detect mismatches and disappearances", WARN),
        ("5", "Modern dashboard", "Dark-themed real-time UI with WebSocket push and historical review", PURPLE),
        ("6", "100% on-premise", "No cloud calls, no rate limits, no subscriptions — pure local intelligence", SUCCESS),
    ]
    for i, (num, title, desc, color) in enumerate(objectives):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.3)
        y = Inches(1.4 + row * 1.85)

        # Number badge
        add_rounded(s, x, y, Inches(0.8), Inches(0.8), color, color, 0)
        add_text(s, x, y, Inches(0.8), Inches(0.8), num, size=24, color=BG, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # Title
        add_text(s, x + Inches(1.0), y, Inches(4.8), Inches(0.45),
                 title, size=18, color=TEXT, bold=True)
        # Description
        add_text(s, x + Inches(1.0), y + Inches(0.55), Inches(4.8), Inches(0.9),
                 desc, size=12, color=TEXT_DIM)

    add_speaker_notes(s, """Walk through each objective. Emphasize that 100% on-premise is not just a feature — it's a privacy stance. Banks, hospitals, and government facilities cannot legally send feeds to AWS. Make sure to mention that ALL six objectives were delivered (not just some).""")

    # ----------------------------------------------------------------
    # SLIDE 5 — System overview
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "System Overview", "From a single RTSP camera to a real-time intelligence dashboard")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 5)

    add_image(s, os.path.join(ASSETS, "workflow.png"),
              Inches(1.5), Inches(1.3), w=Inches(10.3), h=Inches(5.6))

    add_speaker_notes(s, """This is the headline diagram. Use it to anchor the rest of the talk. Walk through the six stages clockwise: capture, detect, recognize, analyze, decide, notify. Stress that all six happen in a single tight loop with sub-second end-to-end latency.""")

    # ----------------------------------------------------------------
    # SLIDE 6 — Architecture
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Clean Architecture", "4 layers · dependency rule · testable · swappable")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 6)

    add_image(s, os.path.join(ASSETS, "architecture.png"),
              Inches(0.5), Inches(1.2), w=Inches(8.5), h=Inches(5.8))

    # Right side panel: why clean arch
    add_rounded(s, Inches(9.2), Inches(1.5), Inches(3.7), Inches(5.5), BG2, GOLD, 2)
    add_text(s, Inches(9.4), Inches(1.7), Inches(3.4), Inches(0.5),
             "Why this matters", size=16, color=GOLD, bold=True)
    add_bullets(s, Inches(9.4), Inches(2.2), Inches(3.4), Inches(4.5), [
        ("Domain layer", "pure entities & interfaces, no I/O"),
        ("Use cases", "orchestrate business rules, framework-agnostic"),
        ("Adapters", "swappable I/O — swap YOLO for any detector"),
        ("Presentation", "FastAPI routes — the only HTTP-aware code"),
        ("Result", "testable, maintainable, extensible"),
    ], size=11, line_spacing=1.2)

    add_speaker_notes(s, """If the committee challenges you on architecture, point out: the domain layer has zero imports from FastAPI, OpenCV, or any framework. You can swap YOLO11s for YOLOv12 or a custom detector by writing a new adapter — the use case doesn't change. This is what makes the project a real engineering effort, not a tutorial copy-paste.""")

    # ----------------------------------------------------------------
    # SLIDE 7 — Tech stack
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Technology Stack", "Battle-tested, production-grade, GPU-accelerated")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 7)

    add_image(s, os.path.join(ASSETS, "tech_stack.png"),
              Inches(0.5), Inches(1.2), w=Inches(12.3), h=Inches(5.8))

    add_speaker_notes(s, """Don't read every item on the slide. Point out: Python 3.11, FastAPI, YOLO11s for primary detection, YOLOv8s weapon model, R3D-18 for violence, InsightFace for face recognition, and ONNX Runtime + CUDA for GPU acceleration. Mention that this combination is rare — most academic projects stop at YOLO + face recognition.""")

    # ----------------------------------------------------------------
    # SLIDE 8 — Object detection deep dive
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Feature 1 · Object Detection", "YOLO11s · 6 classes · per-class confidence")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 8)

    add_image(s, os.path.join(ASSETS, "detection_pipeline.png"),
              Inches(0.5), Inches(1.2), w=Inches(7.5), h=Inches(5.6))

    # Right: details
    add_rounded(s, Inches(8.3), Inches(1.4), Inches(4.6), Inches(2.0), BG2, GOLD, 2)
    add_text(s, Inches(8.5), Inches(1.5), Inches(4.4), Inches(0.4),
             "Why YOLO11s?", size=14, color=GOLD, bold=True)
    add_bullets(s, Inches(8.5), Inches(1.9), Inches(4.4), Inches(1.4), [
        ("18.4 MB model", "lightweight"),
        ("1280 × 1280 input", "small-object accuracy"),
        ("25 ms / frame", "40 FPS capability"),
        ("6 classes", "people, phones, bags, weapons"),
    ], size=10, line_spacing=1.1)

    add_rounded(s, Inches(8.3), Inches(3.6), Inches(4.6), Inches(2.0), BG2, GOLD, 2)
    add_text(s, Inches(8.5), Inches(3.7), Inches(4.4), Inches(0.4),
             "Per-class confidence", size=14, color=GOLD, bold=True)
    add_bullets(s, Inches(8.5), Inches(4.1), Inches(4.4), Inches(1.4), [
        ("Phone: 0.20", "tiny, often blurred"),
        ("Bag: 0.35", "common but distinct"),
        ("Person: 0.45", "balance precision / recall"),
        ("Weapon: 0.50", "rare → high precision"),
    ], size=10, line_spacing=1.1)

    add_rounded(s, Inches(8.3), Inches(5.8), Inches(4.6), Inches(1.2), GOLD, GOLD, 2)
    add_text(s, Inches(8.5), Inches(5.85), Inches(4.4), Inches(1.0),
             "Result: ~25 ms per frame · zero false positives in 2-hour test",
             size=11, color=BG, bold=True, italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_speaker_notes(s, """The interesting design choice is per-class confidence. The phone class needs a lower threshold (0.20) because phones are small, often blurred, and partially occluded. The weapon class needs a higher threshold (0.50) because false positives are worse than false negatives — a wrongly flagged knife causes panic. Show that you understand the trade-offs.""")

    # ----------------------------------------------------------------
    # SLIDE 9 — Face recognition deep dive
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Feature 2 · Face Recognition", "InsightFace · ArcFace-R100 · 512-D embeddings")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 9)

    add_image(s, os.path.join(ASSETS, "face_pipeline.png"),
              Inches(0.5), Inches(1.2), w=Inches(7.5), h=Inches(5.6))

    add_rounded(s, Inches(8.3), Inches(1.4), Inches(4.6), Inches(2.0), BG2, GOLD, 2)
    add_text(s, Inches(8.5), Inches(1.5), Inches(4.4), Inches(0.4),
             "Why InsightFace?", size=14, color=GOLD, bold=True)
    add_bullets(s, Inches(8.5), Inches(1.9), Inches(4.4), Inches(1.4), [
        ("ArcFace-R100", "state-of-the-art accuracy"),
        ("512-D embedding", "compact, fast to compare"),
        ("ONNX + CUDA", "GPU-accelerated, no TF"),
        ("Buffalo-L pack", "RetinaFace detector included"),
    ], size=10, line_spacing=1.1)

    add_rounded(s, Inches(8.3), Inches(3.6), Inches(4.6), Inches(2.0), BG2, GOLD, 2)
    add_text(s, Inches(8.5), Inches(3.7), Inches(4.4), Inches(0.4),
             "Multi-embedding fingerprint", size=14, color=GOLD, bold=True)
    add_bullets(s, Inches(8.5), Inches(4.1), Inches(4.4), Inches(1.4), [
        ("Up to 8 per person", "different angles, lighting"),
        ("Quality-gated", "blur/size filter before save"),
        ("Time-diverse", "spread over multiple visits"),
        ("Cosine distance", "best match across fingerprint"),
    ], size=10, line_spacing=1.1)

    add_rounded(s, Inches(8.3), Inches(5.8), Inches(4.6), Inches(1.2), SUCCESS, SUCCESS, 2)
    add_text(s, Inches(8.5), Inches(5.85), Inches(4.4), Inches(1.0),
             "Auto-register unknowns with 5-s cooldown · promotes to known with one click",
             size=11, color=BG, bold=True, italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_speaker_notes(s, """The key innovation here is the multi-embedding fingerprint. A single embedding is fragile — a different angle breaks it. By storing 8 angle-diverse, lighting-diverse, time-diverse embeddings, we get robust re-identification. Even if the person returns in a year with a different haircut, we still match. This is what most simple face-recognition tutorials miss.""")

    # ----------------------------------------------------------------
    # SLIDE 10 — Threat detection deep dive
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Feature 3 · Threat Detection", "3 models · fusion · severity-graded alerts")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 10)

    add_image(s, os.path.join(ASSETS, "threat_pipeline.png"),
              Inches(0.3), Inches(1.2), w=Inches(9.0), h=Inches(5.7))

    # Right: levels + signals
    add_text(s, Inches(9.5), Inches(1.3), Inches(3.5), Inches(0.4),
             "4 severity levels", size=14, color=GOLD, bold=True)
    levels = [
        ("LOW", "0.20", TEXT_DIM),
        ("MEDIUM", "0.45", INFO),
        ("HIGH", "0.70", WARN),
        ("CRITICAL", "0.95", DANGER),
    ]
    for i, (name, score, color) in enumerate(levels):
        y = Inches(1.7 + i * 0.45)
        add_rounded(s, Inches(9.5), y, Inches(3.5), Inches(0.4), color, color, 0)
        add_text(s, Inches(9.6), y, Inches(2.0), Inches(0.4),
                 name, size=11, color=TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(11.6), y, Inches(1.3), Inches(0.4),
                 f"score ≥ {score}", size=10, color=TEXT, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

    add_text(s, Inches(9.5), Inches(3.7), Inches(3.5), Inches(0.4),
             "Escalation bonus", size=14, color=GOLD, bold=True)
    add_text(s, Inches(9.5), Inches(4.1), Inches(3.5), Inches(2.7),
             "When weapon AND violence co-occur:\n\nscore += 0.20\n\n→ CRITICAL even if neither signal alone crosses the threshold.\n\nThis is what catches a person brandishing a weapon during a fight.",
             size=10, color=TEXT_DIM, italic=True)

    add_speaker_notes(s, """This is the most sophisticated feature. Three different ML models run in parallel on the same frame, and a rule engine fuses their outputs into a single severity score. The escalation bonus is the key insight: a single signal might be ambiguous, but two corroborating signals in the same frame is high confidence. Mention 17 Kinetics-400 violence classes (punch, kick, fight, etc.).""")

    # ----------------------------------------------------------------
    # SLIDE 11 — Threat models detail
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Threat Detection · Model Zoo", "Each model specializes in a different signal")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 11)

    models = [
        ("Weapon Detection", "YOLOv8s · gun.pt", "156 MB", "Fine-tuned Shantanukadam model\nSingle 'weapon' class\n640 px · conf 0.30 → MEDIUM, 0.60 → HIGH", DANGER, 1.0, 1.4),
        ("Pose Estimation", "YOLOv8n-Pose", "6.8 MB", "17 COCO keypoints\nDetects fighting stances\nFists raised, grappling, kicking", INFO, 5.0, 1.4),
        ("Action Recognition", "R3D-18", "73 MB", "3D ResNet · Kinetics-400\n16-frame sliding window\n17 violent classes (punch, kick...)", PURPLE, 9.0, 1.4),
        ("Rule Engine", "Custom Python", "<1 MB", "Score fusion + escalation\n5-s cooldown per (type, level, region)\nSnapshot saved per event", SUCCESS, 1.0, 4.3),
        ("WebSocket Push", "FastAPI WS", "—", "Real-time event to dashboard\nAuto-reconnect on disconnect\nJSON payload per alert", INFO, 5.0, 4.3),
        ("Snapshot Storage", "data/alerts/YYYY-MM-DD/", "—", "JPEG per event\nForensic review\n90% JPEG quality", WARN, 9.0, 4.3),
    ]
    for title, model, size, desc, color, x, y in models:
        add_rounded(s, Inches(x), Inches(y), Inches(3.5), Inches(2.7), BG2, color, 2)
        # Color bar at top
        add_rect(s, Inches(x), Inches(y), Inches(3.5), Inches(0.3), color)
        add_text(s, Inches(x+0.1), Inches(y), Inches(3.3), Inches(0.3),
                 title, size=12, color=TEXT, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(x+0.1), Inches(y+0.35), Inches(3.3), Inches(0.3),
                 model, size=10, color=GOLD, italic=True)
        add_text(s, Inches(x+0.1), Inches(y+0.6), Inches(3.3), Inches(0.3),
                 f"📦 {size}", size=9, color=TEXT_DIM)
        add_text(s, Inches(x+0.1), Inches(y+0.95), Inches(3.3), Inches(1.6),
                 desc, size=10, color=TEXT_DIM)

    add_speaker_notes(s, """Detail the models briefly. Weapon model is from a Kaggle researcher (Shantanukadam). Pose is ultralytics' standard. R3D-18 is a 3D ResNet — chosen because X3D wasn't available in torchvision. The rule engine is the 'glue' that turns model outputs into decisions. WebSocket and snapshot storage are part of the same feature — a threat is only useful if the operator sees it and can prove it happened.""")

    # ----------------------------------------------------------------
    # SLIDE 12 — Theft detection deep dive
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Feature 4 · Theft Detection", "IoU tracker · ownership state machine · forensic alerts")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 12)

    add_image(s, os.path.join(ASSETS, "state_machine.png"),
              Inches(0.3), Inches(1.2), w=Inches(8.5), h=Inches(5.7))

    # Right: rules + features
    add_text(s, Inches(9.0), Inches(1.3), Inches(4.0), Inches(0.4),
             "5 ownership rules", size=14, color=GOLD, bold=True)
    rules = [
        ("CLAIM", "held ≥ 1.5 s → owner", SUCCESS),
        ("DROP", "owner leaves → STATIONARY", INFO),
        ("THEFT", "unknown + missing ≥ 3 s", DANGER),
        ("THEFT", "stranger + missing ≥ 3 s", DANGER),
        ("ABANDON", "idle ≥ 30 s", WARN),
    ]
    for i, (name, desc, color) in enumerate(rules):
        y = Inches(1.7 + i * 0.4)
        add_rounded(s, Inches(9.0), y, Inches(1.0), Inches(0.35), color, color, 0)
        add_text(s, Inches(9.0), y, Inches(1.0), Inches(0.35),
                 name, size=10, color=TEXT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(10.1), y, Inches(2.9), Inches(0.35),
                 desc, size=10, color=TEXT_DIM, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(9.0), Inches(3.9), Inches(4.0), Inches(0.4),
             "Persistence", size=14, color=GOLD, bold=True)
    add_bullets(s, Inches(9.0), Inches(4.3), Inches(4.0), Inches(1.5), [
        "items.json (debounced 1 s)",
        "items_log.json (500 events)",
        "survives restarts",
    ], size=10, line_spacing=1.2)

    add_text(s, Inches(9.0), Inches(5.6), Inches(4.0), Inches(0.4),
             "Mirrored to alerts tab", size=14, color=GOLD, bold=True)
    add_text(s, Inches(9.0), Inches(6.0), Inches(4.0), Inches(0.9),
             "Every THEFT event is also broadcast to the threat dashboard at HIGH level — operators see it in one place.",
             size=10, color=TEXT_DIM, italic=True)

    add_speaker_notes(s, """This is the most novel feature in the project. The state machine is a classic AI pattern (Harel statecharts), but the 5 ownership rules and 5-second event cooldown were carefully tuned to balance false positives (crying wolf) and false negatives (missing real thefts). The mirroring to the alerts tab is a UX detail: operators only check one place.""")

    # ----------------------------------------------------------------
    # SLIDE 13 — Dashboard mockup
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Live Dashboard", "Dark theme · real-time updates · modular UI")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 13)

    add_image(s, os.path.join(ASSETS, "dashboard_mockup.png"),
              Inches(0.5), Inches(1.2), w=Inches(9.5), h=Inches(5.7))

    add_rounded(s, Inches(10.3), Inches(1.4), Inches(2.7), Inches(2.4), BG2, GOLD, 2)
    add_text(s, Inches(10.4), Inches(1.5), Inches(2.5), Inches(0.4),
             "4 Tabs", size=14, color=GOLD, bold=True)
    add_bullets(s, Inches(10.4), Inches(1.9), Inches(2.5), Inches(1.8), [
        "Dashboard",
        "Faces",
        "Alerts",
        "Items",
    ], size=11, line_spacing=1.1)

    add_rounded(s, Inches(10.3), Inches(4.0), Inches(2.7), Inches(2.9), BG2, GOLD, 2)
    add_text(s, Inches(10.4), Inches(4.1), Inches(2.5), Inches(0.4),
             "Real-time", size=14, color=GOLD, bold=True)
    add_bullets(s, Inches(10.4), Inches(4.5), Inches(2.5), Inches(2.3), [
        "MJPEG stream",
        "WebSocket push",
        "Toast notifications",
        "Pulse animation",
        "Auto-reconnect",
    ], size=10, line_spacing=1.2)

    add_speaker_notes(s, """Show the dashboard mockup. The 4 tabs map to the 4 features: Dashboard = everything live, Faces = face registry, Alerts = threats, Items = theft. Mention that the UI is built with vanilla JS — no React, no Vue — because the project is server-rendered (Jinja2) and the JS is small enough to be hand-written. This is an engineering choice, not a limitation.""")

    # ----------------------------------------------------------------
    # SLIDE 14 — Performance
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Performance", "Measured on RTX 3060 Laptop GPU · 1080p input")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 14)

    add_image(s, os.path.join(ASSETS, "performance.png"),
              Inches(0.5), Inches(1.2), w=Inches(8.5), h=Inches(5.7))

    # Right: highlights
    add_rounded(s, Inches(9.2), Inches(1.4), Inches(3.7), Inches(1.5), BG2, SUCCESS, 2)
    add_text(s, Inches(9.3), Inches(1.5), Inches(3.5), Inches(0.4),
             "End-to-end", size=12, color=GOLD, bold=True)
    add_text(s, Inches(9.3), Inches(1.9), Inches(3.5), Inches(1.0),
             "~109 ms per frame\n→ 9 FPS\nAcceptable for surveillance",
             size=11, color=TEXT)

    add_rounded(s, Inches(9.2), Inches(3.0), Inches(3.7), Inches(1.5), BG2, INFO, 2)
    add_text(s, Inches(9.3), Inches(3.1), Inches(3.5), Inches(0.4),
             "Memory", size=12, color=GOLD, bold=True)
    add_text(s, Inches(9.3), Inches(3.5), Inches(3.5), Inches(1.0),
             "~1.8 GB active RAM\nFits any modern PC",
             size=11, color=TEXT)

    add_rounded(s, Inches(9.2), Inches(4.6), Inches(3.7), Inches(1.5), BG2, WARN, 2)
    add_text(s, Inches(9.3), Inches(4.7), Inches(3.5), Inches(0.4),
             "Disk", size=12, color=GOLD, bold=True)
    add_text(s, Inches(9.3), Inches(5.1), Inches(3.5), Inches(1.0),
             "~530 MB models\nSnapshots per event",
             size=11, color=TEXT)

    add_rounded(s, Inches(9.2), Inches(6.2), Inches(3.7), Inches(0.7), GOLD, GOLD, 2)
    add_text(s, Inches(9.3), Inches(6.2), Inches(3.5), Inches(0.7),
             "CPU fallback ready", size=11, color=BG, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_speaker_notes(s, """The numbers tell a story: a single mid-range laptop GPU can do real-time multi-model surveillance. CPU fallback works but is ~6x slower (~1.5 FPS). The big chunks are YOLO (~50 ms total for all 3 models) and R3D-18 (~35 ms for the 16-frame window). If we wanted to push to 30 FPS, we could use TensorRT or replace R3D-18 with a lighter model.""")

    # ----------------------------------------------------------------
    # SLIDE 15 — Challenges & solutions
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Challenges & Solutions", "Engineering problems we encountered and fixed")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 15)

    add_image(s, os.path.join(ASSETS, "challenges.png"),
              Inches(0.3), Inches(1.2), w=Inches(12.7), h=Inches(5.7))

    add_speaker_notes(s, """This slide is the most important for the committee. It shows that we didn't just copy a tutorial — we hit real problems and debugged them. The IoU tracker bug (Python truthy/falsy on classes with __len__) is a great teaching moment. The candidate_taker_since overwrite bug was a logic error we caught with print debugging. Mention that the camera being unreachable in the dev environment is a common gotcha that the server handles gracefully.""")

    # ----------------------------------------------------------------
    # SLIDE 16 — Privacy & security
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Privacy & Security", "Designed for sensitive environments")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 16)

    add_image(s, os.path.join(ASSETS, "privacy.png"),
              Inches(0.3), Inches(1.2), w=Inches(12.7), h=Inches(5.7))

    add_speaker_notes(s, """Emphasize the local-first design. In an era of mass surveillance concerns and GDPR-like regulations, the fact that NOTHING leaves the building is a strong selling point. The audit trail (JSON + snapshots) means every detection is reproducible and reviewable. The fail-safe (camera disconnect doesn't crash the server) means a bad cable doesn't kill the whole system.""")

    # ----------------------------------------------------------------
    # SLIDE 17 — Use cases / impact
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Where This Matters", "From a single shop to nationwide infrastructure")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 17)

    add_image(s, os.path.join(ASSETS, "use_cases.png"),
              Inches(0.3), Inches(1.2), w=Inches(12.7), h=Inches(5.7))

    add_speaker_notes(s, """Connect to the original problem. For each use case, think of a real incident: malls (lost children, shoplifting), banks (vault robbery), airports (security checkpoints), schools (active shooter prevention), warehouses (employee theft accounts for 30% of inventory loss), stadiums (crowd incidents). The system is intentionally general-purpose — a single platform that scales to many verticals.""")

    # ----------------------------------------------------------------
    # SLIDE 18 — Impact stats
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "By The Numbers", "What we built, quantified")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 18)

    add_image(s, os.path.join(ASSETS, "impact_stats.png"),
              Inches(0.3), Inches(1.2), w=Inches(12.7), h=Inches(5.7))

    add_speaker_notes(s, """Use these numbers as ammunition. The committee might ask 'how big is this really?' — 5 AI models, 530 MB on disk, 1.8 GB RAM, 200 events in memory. The fact that it runs on a laptop GPU means the deployment cost per site is essentially zero (most cameras ship with an embedded PC).""")

    # ----------------------------------------------------------------
    # SLIDE 19 — Community impact
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Community Impact", "Real problems this project helps solve")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 19)

    impacts = [
        ("🛡", "Public Safety", "Real-time weapon & violence detection at schools, malls, transit hubs — pre-emptive alerts save lives.", DANGER),
        ("💰", "Economic Loss", "Retail theft costs the global economy $100B+ annually. Item-level tracking recovers inventory and deters thieves.", WARN),
        ("👨‍👩‍👧", "Family Safety", "Find lost children, identify elders with dementia, recognize family members entering restricted zones.", SUCCESS),
        ("⚖", "Justice & Evidence", "Forensic snapshot per event, tamper-evident timestamps, full event log — court-admissible evidence.", INFO),
        ("🌍", "Local-First", "100% on-premise — works in regions with no cloud connectivity. Critical for defense, government, rural sites.", PURPLE),
        ("📚", "Open Research", "Clean code, reproducible setup, well-documented. Other researchers can build on this work for free.", GOLD),
    ]
    for i, (icon, title, desc, color) in enumerate(impacts):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.3)
        y = Inches(1.4 + row * 1.85)
        add_rounded(s, x, y, Inches(5.8), Inches(1.6), BG2, color, 2)
        add_text(s, Inches(x+0.1), y+0.1, Inches(0.8), Inches(0.8),
                 icon, size=28, color=color, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(x+1.0), y+0.1, Inches(4.7), Inches(0.4),
                 title, size=15, color=GOLD, bold=True)
        add_text(s, Inches(x+1.0), y+0.55, Inches(4.7), Inches(1.0),
                 desc, size=10, color=TEXT_DIM)

    add_speaker_notes(s, """This is the 'why does it matter' slide. Don't rush through it. Each of these is a real market with real customers. The committee wants to see that you understand the social context, not just the code. Mention specific statistics where you can: $100B retail theft, etc.""")

    # ----------------------------------------------------------------
    # SLIDE 20 — Future work
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Future Work", "What comes after this graduation project")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 20)

    futures = [
        ("Multi-camera", "Stream from N cameras in parallel · cross-camera tracking with re-identification", INFO),
        ("Edge deployment", "Jetson Orin Nano · Coral TPU · 5 W power budget · 30 FPS sustained", SUCCESS),
        ("Active learning", "Operator corrections feed back into model retraining · weekly improvement cycle", PURPLE),
        ("Mobile alerts", "Telegram / WhatsApp push for HIGH and CRITICAL events · silent escalation", WARN),
        ("License plate", "Add ALPR model · correlate with vehicle detection · stolen-vehicle alerts", DANGER),
        ("Audio analysis", "Scream / glass-break / gunshot sound classification · multi-modal fusion", GOLD),
    ]
    for i, (title, desc, color) in enumerate(futures):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.3)
        y = Inches(1.4 + row * 1.85)
        add_rounded(s, x, y, Inches(5.8), Inches(1.6), BG2, color, 2)
        # Color dot
        add_rounded(s, x+0.2, y+0.2, Inches(0.4), Inches(0.4), color, color, 0)
        add_text(s, x+0.2, y+0.2, Inches(0.4), Inches(0.4),
                 str(i+1), size=14, color=TEXT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x+0.8, y+0.15, Inches(4.8), Inches(0.4),
                 title, size=15, color=TEXT, bold=True)
        add_text(s, x+0.8, y+0.6, Inches(4.8), Inches(0.9),
                 desc, size=10, color=TEXT_DIM)

    add_speaker_notes(s, """Future work shows the committee that you understand this is a starting point. The most realistic near-term wins: Jetson deployment (5W power, fanless) and multi-camera. The active learning loop is a research direction. The audio analysis would unlock a new market (home security). Don't promise to do them all — pick one or two and explain how.""")

    # ----------------------------------------------------------------
    # SLIDE 21 — Lessons learned
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Lessons Learned", "What this project taught us")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 21)

    lessons = [
        ("Clean architecture pays off", "Swapping YOLO11s for a custom model took 1 hour, not 1 day.", SUCCESS),
        ("Bugs hide in defaults", "Python truthy/falsy + classes with __len__ is a real footgun.", WARN),
        ("State machines are the right tool", "Explicit states + transition rules beat implicit flags.", INFO),
        ("WebSocket > polling", "A single bidirectional connection beats 10 polling clients.", PURPLE),
        ("Cooldowns are essential", "Without 5-s cooldown, threat spam would bury the dashboard.", DANGER),
        ("Persistence = trust", "JSON logs + snapshots make the system auditable, not just reactive.", GOLD),
    ]
    for i, lesson in enumerate(lessons):
        title, desc, _ = lesson
        y = Inches(1.5 + i * 0.85)
        # Number
        add_text(s, Inches(0.8), y, Inches(0.6), Inches(0.6),
                 f"{i+1:02d}", size=20, color=GOLD, bold=True)
        # Title
        add_text(s, Inches(1.5), y, Inches(11.5), Inches(0.4),
                 title, size=15, color=TEXT, bold=True)
        # Description
        add_text(s, Inches(1.5), y+0.4, Inches(11.5), Inches(0.4),
                 desc, size=11, color=TEXT_DIM, italic=True)
        if i < len(lessons) - 1:
            add_rect(s, Inches(1.5), y+0.85, Inches(11), Inches(0.01), NAVY)

    add_speaker_notes(s, """Personal reflection goes a long way with the committee. These are not generic platitudes — they came from real moments in the project. The 'clean architecture pays off' is a defense of your design choices. The 'bugs hide in defaults' shows you've debugged. The 'WebSocket > polling' shows architectural maturity.""")

    # ----------------------------------------------------------------
    # SLIDE 22 — Acknowledgements
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Acknowledgements", "Built on the shoulders of giants")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 22)

    add_text(s, Inches(1.0), Inches(1.5), Inches(11.3), Inches(0.5),
             "Open-source models and libraries", size=18, color=GOLD, bold=True)

    sources = [
        ("YOLO11s · YOLOv8s · YOLOv8n-Pose", "Ultralytics — state-of-the-art real-time detection"),
        ("ArcFace-R100 (Buffalo-L)", "InsightFace — face recognition SOTA"),
        ("R3D-18", "Facebook AI Research (Kinetics-400)"),
        ("gun.pt", "Shantanukadam — fine-tuned weapon model"),
        ("FastAPI · Uvicorn · Pydantic", "Modern Python web stack"),
        ("PyTorch · ONNX Runtime · CUDA", "Deep learning inference engine"),
    ]
    for i, (name, desc) in enumerate(sources):
        y = Inches(2.1 + i * 0.55)
        add_text(s, Inches(1.0), y, Inches(5.0), Inches(0.4),
                 name, size=14, color=TEXT, bold=True)
        add_text(s, Inches(6.2), y, Inches(6.6), Inches(0.4),
                 desc, size=12, color=TEXT_DIM, italic=True)
        if i < len(sources) - 1:
            add_rect(s, Inches(1.0), y+0.5, Inches(11.3), Inches(0.01), NAVY)

    add_rounded(s, Inches(1.0), Inches(5.7), Inches(11.3), Inches(1.3), NAVY, GOLD, 1)
    add_text(s, Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.5),
             "Special thanks", size=14, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.0), Inches(6.25), Inches(11.3), Inches(0.5),
             "to our supervisor Dr. [Supervisor Name] for the guidance,\nand to the open-source community for making this work possible.",
             size=12, color=TEXT_DIM, italic=True, align=PP_ALIGN.CENTER)

    add_speaker_notes(s, """Acknowledge that you stood on the shoulders of giants. Name the specific models and their creators. Thank your supervisor. This shows professional humility, which the committee values.""")

    # ----------------------------------------------------------------
    # SLIDE 23 — Conclusion
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Conclusion", "What we delivered")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 23)

    # Big success indicators
    points = [
        ("✓", "Detection", "Real-time 6-class detection on GPU"),
        ("✓", "Recognition", "Multi-embedding face recognition with auto-register"),
        ("✓", "Threat", "3-model fusion with severity scoring and snapshots"),
        ("✓", "Theft", "State-machine ownership tracking with forensic alerts"),
        ("✓", "Dashboard", "Modern dark UI with WebSocket real-time push"),
        ("✓", "Architecture", "Clean 4-layer design, fully testable and swappable"),
    ]
    for i, (check, title, desc) in enumerate(points):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.2)
        y = Inches(1.5 + row * 1.6)
        add_rounded(s, x, y, Inches(3.9), Inches(1.3), BG2, SUCCESS, 2)
        add_text(s, x+0.1, y+0.1, Inches(0.6), Inches(0.6),
                 check, size=28, color=SUCCESS, bold=True)
        add_text(s, x+0.7, y+0.15, Inches(3.1), Inches(0.4),
                 title, size=15, color=GOLD, bold=True)
        add_text(s, x+0.7, y+0.55, Inches(3.1), Inches(0.6),
                 desc, size=10, color=TEXT_DIM)

    # Big takeaway
    add_rounded(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.7), NAVY, GOLD, 2)
    add_text(s, Inches(0.6), Inches(5.15), Inches(12.1), Inches(0.5),
             "Sentinel is a complete, production-grade, AI-powered surveillance system",
             size=18, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(5.65), Inches(12.1), Inches(0.5),
             "that runs on a single laptop, deploys to any site in minutes, and respects privacy by design.",
             size=13, color=TEXT, italic=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.5),
             "Ready for real-world deployment today.",
             size=14, color=SUCCESS, bold=True, italic=True, align=PP_ALIGN.CENTER)

    add_speaker_notes(s, """This is the summary slide. Recap the 6 deliverables. Land the final line with confidence: 'Ready for real-world deployment today.' This is your closing argument. Pause before transitioning to Q&A.""")

    # ----------------------------------------------------------------
    # SLIDE 24 — Thank you / Q&A
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s, BG)
    add_image(s, os.path.join(ASSETS, "title_bg.png"),
              Inches(0), Inches(0), w=SLIDE_W, h=SLIDE_H)
    add_image(s, os.path.join(ASSETS, "logo.png"),
              Inches(5.665), Inches(1.0), w=Inches(2.0), h=Inches(2.0))

    add_text(s, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.9),
             "Thank You", size=72, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
             "Questions & Discussion", size=28, color=TEXT, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(3.5), Inches(5.4), Inches(6.3), Inches(0.04), GOLD)
    add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4),
             "Sentinel · AI-Powered Smart Surveillance System",
             size=16, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.4),
             "Team: [Your Name]    Supervisor: [Dr. Supervisor]",
             size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
             "[Faculty · Department · 2025]",
             size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER, italic=True)

    add_speaker_notes(s, """Thank the committee for their time. Invite questions. Have your demo ready — if they ask 'can you show it', a live demo impresses more than any slide. Be prepared to answer questions on architecture choices, model selection, performance trade-offs, and future work.""")

    # ----------------------------------------------------------------
    # SLIDE 25 — Backup: detailed architecture
    # ----------------------------------------------------------------
    s = prs.slides.add_slide(blank_layout)
    add_bg(s)
    add_gold_accent_top(s)
    add_slide_title(s, "Backup · Detailed Architecture", "For technical questions")
    add_gold_accent_bottom(s)
    add_slide_footer(s, 25)

    add_text(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.5),
             "Code structure — Clean Architecture in practice", size=18, color=GOLD, bold=True)

    code_lines = [
        ("src/", TEXT),
        ("├── domain/", TEXT),
        ("│   ├── entities/        # Camera, Detection, Face, Threat, Item (pure data)", TEXT_DIM),
        ("│   └── interfaces/      # StreamRepo, DetectionRepo, FaceRepo (abstract)", TEXT_DIM),
        ("├── usecases/             # Orchestrate business rules", TEXT),
        ("│   ├── camera_stream.py # RTSP → frame", TEXT_DIM),
        ("│   ├── detection.py     # YOLO inference", TEXT_DIM),
        ("│   ├── face.py          # Recognition + auto-register", TEXT_DIM),
        ("│   ├── threat.py        # Multi-model fusion", TEXT_DIM),
        ("│   └── item_tracking.py # State machine", TEXT_DIM),
        ("├── adapters/", TEXT),
        ("│   ├── rtsp/            # OpenCV reader", TEXT_DIM),
        ("│   ├── yolo/            # 3 YOLO models", TEXT_DIM),
        ("│   ├── insightface/     # Face recognition", TEXT_DIM),
        ("│   ├── torch/           # R3D-18", TEXT_DIM),
        ("│   ├── tracking/        # IoU + ownership", TEXT_DIM),
        ("│   └── storage/         # JSON persistence", TEXT_DIM),
        ("├── presentation/", TEXT),
        ("│   ├── api/             # FastAPI factory", TEXT_DIM),
        ("│   ├── routes/          # 5 routers", TEXT_DIM),
        ("│   ├── realtime/        # WebSocket manager", TEXT_DIM),
        ("│   ├── templates/       # Jinja2", TEXT_DIM),
        ("│   └── static/          # CSS + JS", TEXT_DIM),
        ("└── config/              # Pydantic settings", TEXT),
    ]
    for i, (line, color) in enumerate(code_lines):
        y = Inches(2.0 + i * 0.22)
        add_text(s, Inches(0.8), y, Inches(11.5), Inches(0.22),
                 line, size=10, color=color, font="Consolas")

    add_speaker_notes(s, """Backup slide for technical questions. If the committee asks 'where is X?' you can point them here. The structure clearly shows the 4 layers with examples of what's in each. This is the slide that proves you didn't just write one big main.py.""")

    # ----------------------------------------------------------------
    # Save
    # ----------------------------------------------------------------
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")
    return OUTPUT


if __name__ == "__main__":
    build()
